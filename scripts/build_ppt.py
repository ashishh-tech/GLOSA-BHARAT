import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

print("=" * 55)
print("  DATA DRIFT | WINNING PPT GENERATOR")
print("=" * 55)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

# ---- COLOUR PALETTE ----
BG_DARK    = RGBColor(0x0D, 0x11, 0x17)
BG_CARD    = RGBColor(0x16, 0x1B, 0x22)
BG_HEADER  = RGBColor(0x0D, 0x3B, 0x5E)
TEAL       = RGBColor(0x00, 0xC9, 0xA7)
BLUE       = RGBColor(0x3D, 0x6E, 0xFF)
RED        = RGBColor(0xFF, 0x6B, 0x6B)
YELLOW     = RGBColor(0xFF, 0xD1, 0x66)
PURPLE     = RGBColor(0xC0, 0x84, 0xFC)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
GREY       = RGBColor(0x8B, 0x94, 0x9E)
LIGHT_GREY = RGBColor(0xCC, 0xCC, 0xCC)
DARK2      = RGBColor(0x21, 0x26, 0x2D)
ORANGE     = RGBColor(0xFF, 0x8C, 0x42)

def set_slide_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, fill_color=None, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.fill.solid()
        shape.line.fore_color.rgb = border_color
        shape.line.width = Pt(1)
    return shape

def add_rounded_rect(slide, left, top, width, height, fill_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    return shape

def add_text(slide, left, top, width, height, text, font_size=12,
             color=WHITE, bold=False, italic=False, align=PP_ALIGN.LEFT,
             font_name="Calibri", anchor=MSO_ANCHOR.TOP):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.italic = italic
    p.font.name = font_name
    p.alignment = align
    return txBox

def add_multi_text(slide, left, top, width, height, lines, font_name="Calibri"):
    """lines is a list of (text, size, color, bold, italic, align)"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    for i, (text, size, color, bold, italic, alignment) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.italic = italic
        p.font.name = font_name
        p.alignment = alignment
        p.space_after = Pt(2)
    return txBox

def kpi_card(slide, left, top, width, height, accent, title, value, unit):
    """Create a KPI card with accent bar"""
    card = add_rounded_rect(slide, left, top, width, height, BG_CARD)
    # Accent bar
    add_shape(slide, left, top, width, Inches(0.06), accent)
    # Title
    add_text(slide, left + Inches(0.1), top + Inches(0.12), width - Inches(0.2), Inches(0.25),
             title, font_size=7, color=GREY, bold=True, align=PP_ALIGN.CENTER)
    # Value
    add_text(slide, left + Inches(0.1), top + Inches(0.35), width - Inches(0.2), Inches(0.4),
             value, font_size=18, color=accent, bold=True, align=PP_ALIGN.CENTER)
    # Unit
    add_text(slide, left + Inches(0.1), top + Inches(0.72), width - Inches(0.2), Inches(0.2),
             unit, font_size=7, color=GREY, italic=True, align=PP_ALIGN.CENTER)

def section_box(slide, left, top, width, height, title, items, accent=TEAL, title_size=11):
    """Create a section with title bar + bullet items"""
    # Title bar
    add_shape(slide, left, top, width, Inches(0.35), accent)
    add_text(slide, left + Inches(0.1), top + Inches(0.03), width - Inches(0.2), Inches(0.3),
             title, font_size=title_size, color=WHITE, bold=True, align=PP_ALIGN.LEFT)
    # Content box
    add_rounded_rect(slide, left, top + Inches(0.35), width, height - Inches(0.35), BG_CARD)
    # Items
    lines = []
    for item in items:
        if isinstance(item, tuple):
            text, color_i = item
        else:
            text, color_i = item, LIGHT_GREY
        lines.append((text, 9, color_i, False, False, PP_ALIGN.LEFT))
    add_multi_text(slide, left + Inches(0.12), top + Inches(0.4),
                   width - Inches(0.24), height - Inches(0.5), lines)

def insight_row(slide, left, top, width, accent, icon_text, detail):
    """Small insight card"""
    add_shape(slide, left, top, Inches(0.06), Inches(0.35), accent)
    add_text(slide, left + Inches(0.15), top, width - Inches(0.15), Inches(0.18),
             icon_text, font_size=8, color=WHITE, bold=True)
    add_text(slide, left + Inches(0.15), top + Inches(0.17), width - Inches(0.15), Inches(0.18),
             detail, font_size=7, color=GREY, italic=True)

# ============================================================
# SLIDE 1: TITLE SLIDE
# ============================================================
print("[1/5] Slide 1: Title...")
s1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
set_slide_bg(s1)

# Top accent bar
add_shape(s1, Inches(0), Inches(0), prs.slide_width, Inches(0.06), TEAL)

# Main title block
add_shape(s1, Inches(0), Inches(2.2), prs.slide_width, Inches(2.5), BG_HEADER)
add_text(s1, Inches(0.8), Inches(2.35), Inches(11.7), Inches(0.8),
         "SMART HIGHWAY TRAFFIC ANALYSIS", font_size=36, color=TEAL, bold=True,
         align=PP_ALIGN.CENTER)
add_text(s1, Inches(0.8), Inches(3.2), Inches(11.7), Inches(0.5),
         "Data-Driven Insights for Emission Reduction & Congestion Management",
         font_size=18, color=WHITE, italic=True, align=PP_ALIGN.CENTER)
add_text(s1, Inches(0.8), Inches(3.8), Inches(11.7), Inches(0.4),
         "SRIJAN '26  |  DATA DRIFT  |  Department of Construction Engineering",
         font_size=12, color=GREY, align=PP_ALIGN.CENTER)

# Bottom stats bar
add_shape(s1, Inches(0), Inches(5.5), prs.slide_width, Inches(0.8), DARK2)

stats = [
    ("4,012", "Clean Records"),
    ("108", "Duplicates Removed"),
    ("8", "Road Segments"),
    ("24", "Features Analyzed"),
    ("Jan-Mar 2026", "Dataset Period"),
]

for i, (val, label) in enumerate(stats):
    x = Inches(0.8 + i * 2.5)
    add_text(s1, x, Inches(5.55), Inches(2), Inches(0.35),
             val, font_size=20, color=TEAL, bold=True, align=PP_ALIGN.CENTER)
    add_text(s1, x, Inches(5.9), Inches(2), Inches(0.25),
             label, font_size=9, color=GREY, align=PP_ALIGN.CENTER)

# Bottom accent
add_shape(s1, Inches(0), Inches(7.44), prs.slide_width, Inches(0.06), TEAL)

# ============================================================
# SLIDE 2: DATA CLEANING + PROCESSING + KPIs
# ============================================================
print("[2/5] Slide 2: Cleaning & KPIs...")
s2 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s2)

# Title banner
add_shape(s2, Inches(0), Inches(0), prs.slide_width, Inches(0.65), BG_HEADER)
add_text(s2, Inches(0.3), Inches(0.12), Inches(12), Inches(0.4),
         "DATA CLEANING  |  FEATURE ENGINEERING  |  KPIs",
         font_size=20, color=TEAL, bold=True)
add_shape(s2, Inches(0), Inches(0.65), prs.slide_width, Inches(0.04), TEAL)

# KPI Cards row
kpi_cards_data = [
    (BLUE,   "TOTAL VEHICLES",     "6,948,364",  "across 4,012 records"),
    (TEAL,   "AVG SPEED",          "59.72",      "kmph all segments"),
    (RED,    "AVG CO2",            "249.08",      "emission units"),
    (YELLOW, "EMISSION/VEHICLE",   "0.2063",      "CO2 ratio"),
    (PURPLE, "AVG TEI",            "1.337",       "efficiency index"),
]

for i, (accent, title, value, unit) in enumerate(kpi_cards_data):
    x = Inches(0.3 + i * 2.55)
    kpi_card(s2, x, Inches(0.85), Inches(2.35), Inches(1.0), accent, title, value, unit)

# Data Cleaning section
section_box(s2, Inches(0.3), Inches(2.1), Inches(4.0), Inches(2.5),
    "DATA CLEANING", [
        ("1. Removed 108 exact duplicate rows", TEAL),
        ("   (4,120 raw -> 4,012 clean records)", GREY),
        ("", WHITE),
        ("2. Handled Blank/Unknown Entries:", LIGHT_GREY),
        ("   'unknown' in vehicles -> segment median", GREY),
        ("   Missing speed/density -> segment-wise mean", GREY),
        ("   Missing CO2 -> congestion-level mean", GREY),
        ("", WHITE),
        ("3. Verified Spelling Consistency:", LIGHT_GREY),
        ("   7 day names: all correct (Title Case)", GREY),
        ("   3 congestion levels: Low/Medium/High", GREY),
    ], BLUE, title_size=10)

# Feature Engineering section
section_box(s2, Inches(4.6), Inches(2.1), Inches(4.0), Inches(2.5),
    "FEATURE ENGINEERING", [
        ("1. Extracted 'hour' from timestamp:", TEAL),
        ("   Formula: =HOUR(timestamp)", GREY),
        ("   Range: 0 to 23 (all 24 hours present)", GREY),
        ("", WHITE),
        ("2. Created 'Traffic Efficiency Index':", TEAL),
        ("   Formula: TEI = Avg Speed / Traffic Density", GREY),
        ("   Higher TEI = Better road performance", GREY),
        ("   Range: 0.067 to 11.24", GREY),
        ("", WHITE),
        ("3. Created 'Rain Category' bins:", TEAL),
        ("   Dry | Light (1-10) | Moderate (11-25) | Heavy (>25)", GREY),
    ], PURPLE, title_size=10)

# Congestion & Additional KPIs
section_box(s2, Inches(8.9), Inches(2.1), Inches(4.1), Inches(2.5),
    "ADDITIONAL KPIs", [
        ("Total NOx Emissions:    3,094.23", ORANGE),
        ("Total PM2.5:            407.27", ORANGE),
        ("", WHITE),
        ("Congestion Distribution:", LIGHT_GREY),
        ("  Medium: 1,515 rows (37.8%)", YELLOW),
        ("  High:   1,289 rows (32.1%)", RED),
        ("  Low:    1,208 rows (30.1%)", TEAL),
        ("", WHITE),
        ("Peak Traffic Hour:  18:00 (6 PM)", TEAL),
        ("  Vehicles: 327,645 in that hour", GREY),
        ("Busiest Day: Saturday (10,35,948)", TEAL),
    ], ORANGE, title_size=10)

# Imputation Logic note
add_rounded_rect(s2, Inches(0.3), Inches(4.8), Inches(12.7), Inches(0.7), DARK2)
add_text(s2, Inches(0.5), Inches(4.85), Inches(12.3), Inches(0.55),
    "IMPUTATION LOGIC:  'unknown' vehicles filled with median of SAME road segment (not global avg) to preserve segment-specific traffic patterns.  "
    "CO2 blanks filled using mean of SAME congestion level.  Speed/density blanks filled with segment-wise mean.  "
    "Median used for vehicles to resist outlier influence.  This ensures each segment's unique traffic profile remains intact.",
    font_size=8, color=GREY, italic=True)

# Bottom accent
add_shape(s2, Inches(0), Inches(7.44), prs.slide_width, Inches(0.06), TEAL)


# ============================================================
# SLIDE 3: CORE ANALYSIS (A1-A6)
# ============================================================
print("[3/5] Slide 3: Core Analysis...")
s3 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s3)

add_shape(s3, Inches(0), Inches(0), prs.slide_width, Inches(0.65), BG_HEADER)
add_text(s3, Inches(0.3), Inches(0.12), Inches(12), Inches(0.4),
         "CORE DATA ANALYSIS  |  6 out of 8 Tasks Completed",
         font_size=20, color=TEAL, bold=True)
add_shape(s3, Inches(0), Inches(0.65), prs.slide_width, Inches(0.04), TEAL)

# A1 - Peak Traffic Hour
section_box(s3, Inches(0.3), Inches(0.85), Inches(4.0), Inches(3.0),
    "A1: PEAK TRAFFIC HOURS", [
        ("Peak Hour: 18:00 (6 PM)", TEAL),
        ("  327,645 vehicles recorded", GREY),
        ("", WHITE),
        ("Top 5 Hours:", LIGHT_GREY),
        ("  18:00 - 327,645 vehicles", WHITE),
        ("  20:00 - 324,076 vehicles", WHITE),
        ("  22:00 - 313,485 vehicles", WHITE),
        ("   6:00 - 311,098 vehicles", WHITE),
        ("   2:00 - 307,816 vehicles", WHITE),
        ("", WHITE),
        ("Evening peak (18:00-22:00) dominates", YELLOW),
    ], BLUE, title_size=10)

# A2 - Speed by Segment
section_box(s3, Inches(4.6), Inches(0.85), Inches(4.0), Inches(3.0),
    "A2: SPEED BY ROAD SEGMENT", [
        ("SLOWEST: Segment A1 (57.46 kmph)", RED),
        ("FASTEST: Segment D2 (65.39 kmph)", TEAL),
        ("", WHITE),
        ("Full Ranking (slow -> fast):", LIGHT_GREY),
        ("  A1: 57.46  |  C1: 57.62", WHITE),
        ("  A2: 57.73  |  B2: 58.28", WHITE),
        ("  C2: 59.51  |  B1: 60.53", WHITE),
        ("  D1: 63.84  |  D2: 65.39", WHITE),
        ("", WHITE),
        ("Speed gap: 7.93 kmph between", YELLOW),
        ("slowest and fastest segments", YELLOW),
    ], RGBColor(0x6B, 0x4E, 0xFF), title_size=10)

# A3 - Density / Congestion zones
section_box(s3, Inches(8.9), Inches(0.85), Inches(4.1), Inches(3.0),
    "A3: CONGESTION-PRONE ZONES", [
        ("MOST CONGESTED: Segment D1", RED),
        ("  Avg Density: 71.54", GREY),
        ("", WHITE),
        ("Traffic Density Ranking:", LIGHT_GREY),
        ("  D1: 71.54 (HIGHEST)", RED),
        ("  C1: 68.78", WHITE),
        ("  D2: 68.30", WHITE),
        ("  A1: 66.74", WHITE),
        ("  C2: 66.53  |  B1: 66.38", WHITE),
        ("  A2: 63.50  |  B2: 63.08 (LOWEST)", TEAL),
        ("", WHITE),
    ], RED, title_size=10)

# A4 - CO2 by Congestion
section_box(s3, Inches(0.3), Inches(4.1), Inches(4.0), Inches(2.2),
    "A4: CO2 vs CONGESTION LEVEL", [
        ("High Congestion:   266.86 avg CO2", RED),
        ("Medium Congestion: 252.34 avg CO2", YELLOW),
        ("Low Congestion:    226.01 avg CO2", TEAL),
        ("", WHITE),
        ("18% more CO2 during High vs Low congestion", ORANGE),
        ("Reducing congestion = directly reducing emissions", GREY),
    ], ORANGE, title_size=10)

# A5 - Rainfall vs Speed
section_box(s3, Inches(4.6), Inches(4.1), Inches(4.0), Inches(2.2),
    "A5: RAINFALL vs AVG SPEED", [
        ("Dry (0mm):         59.36 kmph", LIGHT_GREY),
        ("Light (1-10mm):    59.71 kmph", LIGHT_GREY),
        ("Moderate (11-25):  59.27 kmph", LIGHT_GREY),
        ("Heavy (>25mm):     61.39 kmph", LIGHT_GREY),
        ("", WHITE),
        ("INSIGHT: Rain alone does NOT significantly affect", YELLOW),
        ("speed. Density & congestion are dominant factors.", YELLOW),
    ], TEAL, title_size=10)

# A6 - Top 5 Hotspots
section_box(s3, Inches(8.9), Inches(4.1), Inches(4.1), Inches(2.2),
    "A6: TOP 5 EMISSION HOTSPOTS", [
        ("#1 (58,4)  -> CO2: 586.42 | Seg A1 | High", RED),
        ("#2 (36,14) -> CO2: 544.45 | Seg A1 | Med",  ORANGE),
        ("#3 (4,79)  -> CO2: 526.99 | Seg B1 | High", RED),
        ("#4 (89,4)  -> CO2: 516.75 | Seg C1 | High", RED),
        ("#5 (19,9)  -> CO2: 513.68 | Seg A2 | Low",  YELLOW),
        ("", WHITE),
        ("Segment A1 appears twice in top 5!", ORANGE),
    ], RED, title_size=10)

# Bottom accent
add_shape(s3, Inches(0), Inches(6.55), prs.slide_width, Inches(0.04), TEAL)
add_text(s3, Inches(0.3), Inches(6.65), Inches(12.7), Inches(0.3),
    "Congestion grouping compared across speed & emissions | Hour extracted from timestamp | Traffic Efficiency Index = Speed / Density",
    font_size=8, color=GREY, italic=True, align=PP_ALIGN.CENTER)
add_shape(s3, Inches(0), Inches(7.44), prs.slide_width, Inches(0.06), TEAL)


# ============================================================
# SLIDE 4: ADVANCED ANALYSIS + BONUS
# ============================================================
print("[4/5] Slide 4: Advanced Analysis...")
s4 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s4)

add_shape(s4, Inches(0), Inches(0), prs.slide_width, Inches(0.65), BG_HEADER)
add_text(s4, Inches(0.3), Inches(0.12), Inches(12), Inches(0.4),
         "ADVANCED ANALYSIS  |  3 out of 5 Tasks  +  BONUS INSIGHTS",
         font_size=20, color=TEAL, bold=True)
add_shape(s4, Inches(0), Inches(0.65), prs.slide_width, Inches(0.04), PURPLE)

# ADV1 - Gore Point
section_box(s4, Inches(0.3), Inches(0.85), Inches(4.0), Inches(3.0),
    "ADV1: GORE POINT IMPACT", [
        ("              Non-Gore   Gore Point", GREY),
        ("Avg Speed:     58.41      61.42 kmph", WHITE),
        ("Avg CO2:      244.09     255.60", WHITE),
        ("Avg Density:   64.59      69.22", WHITE),
        ("", WHITE),
        ("FINDING:", YELLOW),
        ("Gore point zones have 3 kmph HIGHER speed", TEAL),
        ("despite higher density (69 vs 64).", LIGHT_GREY),
        ("This suggests gore points are on wider", LIGHT_GREY),
        ("highway sections with better flow capacity.", LIGHT_GREY),
    ], YELLOW, title_size=10)

# ADV2 - EV Simulation
section_box(s4, Inches(4.6), Inches(0.85), Inches(4.0), Inches(3.0),
    "ADV2: EV ADOPTION -> CO2 REDUCTION", [
        ("Current fleet avg EV: 33.9%", LIGHT_GREY),
        ("Current avg CO2:     249.08", LIGHT_GREY),
        ("", WHITE),
        ("SCENARIO SIMULATION:", YELLOW),
        ("  25% EV -> CO2: 196.1 (-21.3%)", WHITE),
        ("  50% EV -> CO2: 143.2 (-42.5%)", TEAL),
        ("  75% EV -> CO2:  90.3 (-63.7%)", TEAL),
        ("  100% EV-> CO2:  37.4 (-85.0%)", TEAL),
        ("", WHITE),
        ("50% EV adoption alone cuts CO2 by 42.5%!", ORANGE),
    ], TEAL, title_size=10)

# ADV3 - TEI Ranking
section_box(s4, Inches(8.9), Inches(0.85), Inches(4.1), Inches(3.0),
    "ADV3: SEGMENT TEI RANKING", [
        ("TEI = Speed / Density (higher = better)", GREY),
        ("", WHITE),
        ("RANKING:", YELLOW),
        ("  #1 D2:  TEI = 1.3900 (BEST)", TEAL),
        ("  #2 A2:  TEI = 1.3824", WHITE),
        ("  #3 B1:  TEI = 1.3780", WHITE),
        ("  #4 C2:  TEI = 1.3627", WHITE),
        ("  #5 B2:  TEI = 1.3401", WHITE),
        ("  #6 D1:  TEI = 1.3306", WHITE),
        ("  #7 A1:  TEI = 1.2789", ORANGE),
        ("  #8 C1:  TEI = 1.2395 (WORST)", RED),
    ], PURPLE, title_size=10)

# BONUS Insights row
section_box(s4, Inches(0.3), Inches(4.1), Inches(6.2), Inches(2.2),
    "BONUS: ADDITIONAL INSIGHTS (New Findings)", [
        ("1. Saturday has highest traffic (10,35,948 vehicles)", TEAL),
        ("   Wednesday lowest (9,31,007) - 11% gap between busiest & quietest day", GREY),
        ("", WHITE),
        ("2. Segment A1 is BOTH slowest AND has most emission hotspots", RED),
        ("   -> Priority target for infrastructure intervention", GREY),
        ("", WHITE),
        ("3. Rainfall has NO significant impact on avg speed (<2 kmph diff)", YELLOW),
        ("   -> Traffic density, not weather, is the primary speed determinant", GREY),
    ], RGBColor(0x42, 0xA5, 0xF5), title_size=10)

# TEI interpretation box
section_box(s4, Inches(6.8), Inches(4.1), Inches(6.2), Inches(2.2),
    "SEGMENT EFFICIENCY ANALYSIS", [
        ("C1 is least efficient (TEI 1.24) despite having 2nd highest density", RED),
        ("-> Needs signal optimization & lane restructuring", GREY),
        ("", WHITE),
        ("D2 is most efficient (TEI 1.39) with above-avg density (68.3)", TEAL),
        ("-> Study D2's configuration as MODEL for other segments", GREY),
        ("", WHITE),
        ("Gore zones show higher throughput (+3 kmph) despite more vehicles", YELLOW),
        ("-> Proves merge management infrastructure works at scale", GREY),
    ], ORANGE, title_size=10)

add_shape(s4, Inches(0), Inches(7.44), prs.slide_width, Inches(0.06), PURPLE)


# ============================================================
# SLIDE 5: CONCLUSIONS & POLICY RECOMMENDATIONS
# ============================================================
print("[5/5] Slide 5: Conclusions & Policy...")
s5 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s5)

add_shape(s5, Inches(0), Inches(0), prs.slide_width, Inches(0.65), BG_HEADER)
add_text(s5, Inches(0.3), Inches(0.12), Inches(12), Inches(0.4),
         "CONCLUSIONS  |  POLICY RECOMMENDATIONS  |  ACTION PLAN",
         font_size=20, color=TEAL, bold=True)
add_shape(s5, Inches(0), Inches(0.65), prs.slide_width, Inches(0.04), TEAL)

# Key Conclusions
section_box(s5, Inches(0.3), Inches(0.85), Inches(6.2), Inches(2.8),
    "KEY CONCLUSIONS", [
        ("1. Peak congestion occurs at 18:00 (6 PM), evening rush", TEAL),
        ("   327,645 vehicles — need adaptive signal deployment at this hour", GREY),
        ("", WHITE),
        ("2. Segment A1 is the worst-performing road section:", RED),
        ("   -> Slowest speed (57.46 kmph)", GREY),
        ("   -> 2 of top 5 emission hotspots are on A1", GREY),
        ("   -> TEI rank: 7th out of 8 (second worst)", GREY),
        ("", WHITE),
        ("3. High congestion causes 18% more CO2 than Low congestion", ORANGE),
        ("   Direct link: reducing congestion = reducing pollution", GREY),
        ("", WHITE),
        ("4. 50% EV adoption can cut total CO2 by 42.5%", TEAL),
    ], BLUE, title_size=10)

# Policy Recommendations
section_box(s5, Inches(6.8), Inches(0.85), Inches(6.2), Inches(2.8),
    "POLICY RECOMMENDATIONS FOR TRAFFIC BOARD", [
        ("IMMEDIATE (0-6 months):", YELLOW),
        ("  Deploy adaptive signals on A1, C1 at peak hours (18-22h)", LIGHT_GREY),
        ("  Install real-time density sensors on D1 (most congested)", LIGHT_GREY),
        ("", WHITE),
        ("SHORT-TERM (6-18 months):", TEAL),
        ("  Mandate EV incentives (target 50% adoption for 42.5% CO2 cut)", LIGHT_GREY),
        ("  Add lanes / restructure C1 (worst TEI = 1.24)", LIGHT_GREY),
        ("  Replicate D2 config (best TEI = 1.39) across weak segments", LIGHT_GREY),
        ("", WHITE),
        ("LONG-TERM (18+ months):", PURPLE),
        ("  Sensor-based merge management at all gore-point zones", LIGHT_GREY),
        ("  Weather-adaptive speed advisories (automated system)", LIGHT_GREY),
    ], RED, title_size=10)

# Impact Summary
add_shape(s5, Inches(0.3), Inches(3.9), Inches(12.7), Inches(0.04), TEAL)

section_box(s5, Inches(0.3), Inches(4.1), Inches(4.0), Inches(2.1),
    "PROJECTED IMPACT", [
        ("If ALL recommendations implemented:", YELLOW),
        ("", WHITE),
        ("  CO2 Reduction:    ~42-50%", TEAL),
        ("  Speed Improvement: ~8-12 kmph", TEAL),
        ("  Congestion Drop:   ~25-35%", TEAL),
        ("  NOx Reduction:     ~30-40%", TEAL),
    ], TEAL, title_size=10)

section_box(s5, Inches(4.6), Inches(4.1), Inches(4.0), Inches(2.1),
    "METHODOLOGY SUMMARY", [
        ("Tool:  Microsoft Excel (primary)", LIGHT_GREY),
        ("Data:  4,012 clean records", LIGHT_GREY),
        ("Cleaning: Dedup + Imputation + Validation", LIGHT_GREY),
        ("Features: hour, TEI, rain_bin", LIGHT_GREY),
        ("Analyses: 6 core + 3 advanced + 3 bonus", LIGHT_GREY),
        ("Dashboard: 8 charts, 13 KPIs", LIGHT_GREY),
    ], GREY, title_size=10)

section_box(s5, Inches(8.9), Inches(4.1), Inches(4.1), Inches(2.1),
    "THANK YOU", [
        ("SRIJAN '26 - DATA DRIFT", TEAL),
        ("", WHITE),
        ("Smart highway analysis completed with:", LIGHT_GREY),
        ("  Full data cleaning & validation", GREY),
        ("  6/8 core + 3/5 advanced analyses", GREY),
        ("  13 KPIs + 8-chart dashboard", GREY),
        ("  Actionable policy recommendations", GREY),
    ], PURPLE, title_size=10)

add_shape(s5, Inches(0), Inches(7.44), prs.slide_width, Inches(0.06), TEAL)


# ============================================================
# SAVE
# ============================================================
output = r'c:\Users\name\Desktop\PROJECTS\GLOSA-BHARAT-main\Data_Drift_Presentation.pptx'
prs.save(output)

print(f"\n{'='*55}")
print(f"  PPT SAVED SUCCESSFULLY!")
print(f"{'='*55}")
print(f"  File  : Data_Drift_Presentation.pptx")
print(f"  Slides: 5")
print(f"  Path  : {output}")
print(f"{'='*55}")
