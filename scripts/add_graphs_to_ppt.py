import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
import seaborn as sns
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import warnings
warnings.filterwarnings('ignore')

print("Generating Dark-Themed Graphs...")

# ---------------------------------------------------------
# 1. LOAD DATA & CALC (Same logic as dashboard)
# ---------------------------------------------------------
df = pd.read_csv(r'c:\Users\name\Desktop\PROJECTS\GLOSA-BHARAT-main\highway dataset.csv').drop_duplicates()
df['vehicles'] = pd.to_numeric(df['vehicles'], errors='coerce')
df['congestion'] = df['congestion'].str.strip().str.title()
df['vehicles'] = df.groupby('road segment id')['vehicles'].transform(lambda x: x.fillna(x.median()))
df['traffic density'] = df.groupby('road segment id')['traffic density'].transform(lambda x: x.fillna(x.median()))
df['avg speed kmph'] = df.groupby('road segment id')['avg speed kmph'].transform(lambda x: x.fillna(x.mean()))
df['CO2 emission'] = df.groupby('congestion')['CO2 emission'].transform(lambda x: x.fillna(x.mean()))
df['timestamp'] = pd.to_datetime(df['timestamp'])
df['hour'] = df['timestamp'].dt.hour
df['tei'] = (df['avg speed kmph'] / df['traffic density'].replace(0, np.nan)).fillna(df['avg speed kmph'] / df['traffic density'].median())

# GLOSA Sim Logic
df['glosa_applied'] = ((df['road segment id'].isin(['A1', 'C1'])) & (df['traffic density'] > 65))
df['glosa_co2'] = df['CO2 emission'].copy()
df.loc[df['glosa_applied'], 'glosa_co2'] *= 0.75

os.makedirs("graphs", exist_ok=True)

# ---------------------------------------------------------
# 2. PLOTTING SETUP (Dark Theme)
# ---------------------------------------------------------
# Match PPT dark background (0x0A, 0x0F, 0x1A -> #0A0F1A)
bg_color = "#0B111D"
text_color = "#FFFFFF"
grid_color = "#1E293E"

plt.rcParams.update({
    "figure.facecolor": bg_color,
    "axes.facecolor": bg_color,
    "savefig.facecolor": bg_color,
    "axes.edgecolor": grid_color,
    "axes.labelcolor": text_color,
    "xtick.color": text_color,
    "ytick.color": text_color,
    "text.color": text_color,
    "axes.grid": True,
    "grid.color": grid_color,
    "grid.alpha": 0.5,
})

def format_ax(ax):
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.spines['left'].set_color(grid_color)
    ax.spines['bottom'].set_color(grid_color)

# ---- Chart 1: Peak Traffic Hours ----
peak_hours = df.groupby('hour')['vehicles'].sum()
fig, ax = plt.subplots(figsize=(7, 4))
colors = ["#457BFF" if h != 18 else "#00D2A8" for h in peak_hours.index]
ax.bar(peak_hours.index, peak_hours.values, color=colors)
ax.set_title("Peak Traffic Volume by Hour", fontsize=14, fontweight='bold', pad=15)
ax.set_xlabel("Hour of Day")
ax.set_ylabel("Total Vehicles")
ax.set_xticks(range(0, 24, 2))
format_ax(ax)
plt.tight_layout()
plt.savefig("graphs/1_peak_hours.png", dpi=200)
plt.close()

# ---- Chart 2: Speed by Segment ----
speed_seg = df.groupby('road segment id')['avg speed kmph'].mean().sort_values()
fig, ax = plt.subplots(figsize=(7, 4))
colors = ["#FF5C5C" if s == "A1" else "#457BFF" for s in speed_seg.index]
ax.bar(speed_seg.index, speed_seg.values, color=colors)
ax.set_title("Average Speed by Road Segment", fontsize=14, fontweight='bold', pad=15)
ax.set_ylabel("Speed (kmph)")
ax.set_ylim(55, 68)
format_ax(ax)
plt.tight_layout()
plt.savefig("graphs/2_speed_seg.png", dpi=200)
plt.close()

# ---- Chart 3: CO2 by Congestion ----
co2_cong = df.groupby('congestion')['CO2 emission'].mean()
co2_cong = co2_cong.reindex(["Low", "Medium", "High"])
fig, ax = plt.subplots(figsize=(7, 4))
ax.bar(co2_cong.index, co2_cong.values, color=["#00E676", "#FFC107", "#FF5C5C"], width=0.5)
ax.set_title("CO2 Emissions by Congestion Level", fontsize=14, fontweight='bold', pad=15)
ax.set_ylabel("Avg CO2")
format_ax(ax)
plt.tight_layout()
plt.savefig("graphs/3_co2_cong.png", dpi=200)
plt.close()

# ---- Chart 4: TEI Ranking ----
tei_rank = df.groupby('road segment id')['tei'].mean().sort_values(ascending=True)
fig, ax = plt.subplots(figsize=(7, 4))
colors = ["#FF5C5C" if s == "C1" else ("#00E676" if s == "D2" else "#BB86FC") for s in tei_rank.index]
ax.barh(tei_rank.index, tei_rank.values, color=colors)
ax.set_title("Traffic Efficiency Index (TEI) Ranking", fontsize=14, fontweight='bold', pad=15)
ax.set_xlabel("TEI (Higher is better)")
ax.set_xlim(1.2, 1.45)
format_ax(ax)
plt.tight_layout()
plt.savefig("graphs/4_tei_rank.png", dpi=200)
plt.close()

# ---- Chart 5: GLOSA Simulation Before/After ----
segs = ["A1", "C1"]
b_co2 = [df[df['road segment id']=='A1']['CO2 emission'].mean(), df[df['road segment id']=='C1']['CO2 emission'].mean()]
a_co2 = [df[df['road segment id']=='A1']['glosa_co2'].mean(), df[df['road segment id']=='C1']['glosa_co2'].mean()]
x = np.arange(len(segs))
width = 0.35

fig, ax = plt.subplots(figsize=(7, 4))
r1 = ax.bar(x - width/2, b_co2, width, label='Before GLOSA', color="#FF5C5C")
r2 = ax.bar(x + width/2, a_co2, width, label='After GLOSA', color="#00D2A8")
ax.set_title("GLOSA Deployment: CO2 Drop", fontsize=14, fontweight='bold', pad=15)
ax.set_xticks(x)
ax.set_xticklabels(segs)
ax.set_ylabel("Avg CO2")
ax.legend(facecolor=bg_color, edgecolor=bg_color, labelcolor=text_color)
format_ax(ax)
plt.tight_layout()
plt.savefig("graphs/5_glosa.png", dpi=200)
plt.close()

# ---- Chart 6: EV Simulation ----
curr_ev = df['ev percentage'].mean()
scenarios = ['Current', '25%', '50%', '75%', '100%']
ev_vals = [curr_ev, 0.25, 0.50, 0.75, 1.00]
co2_base = df['CO2 emission'].mean()
co2_sims = [co2_base*(1 - e*0.85) for e in ev_vals]

fig, ax = plt.subplots(figsize=(7, 4))
ax.plot(scenarios, co2_sims, marker='o', linewidth=3, color="#00D2A8", markersize=10)
ax.set_title("EV Adoption Impact on CO2", fontsize=14, fontweight='bold', pad=15)
ax.set_ylabel("Total Avg CO2")
format_ax(ax)
plt.tight_layout()
plt.savefig("graphs/6_ev_sim.png", dpi=200)
plt.close()

print("Graphs created successfully!")

# ---------------------------------------------------------
# 3. APPEND TO PPT
# ---------------------------------------------------------
print("Appending graphs to Presentation...")
pptx_src = r'c:\Users\name\Desktop\PROJECTS\GLOSA-BHARAT-main\Final_Presentation.pptx'
pptx_out = r'c:\Users\name\Desktop\PROJECTS\GLOSA-BHARAT-main\Final_Presentation_With_Graphs.pptx'

prs = Presentation(pptx_src)

# Define PPT Colors 
BG_C = RGBColor(0x0A, 0x0F, 0x1A)
BLUE_C = RGBColor(0x45, 0x7B, 0xFF)
ACCENT_C = RGBColor(0x00, 0xD2, 0xA8)
WHT_C = RGBColor(0xFF, 0xFF, 0xFF)

def add_graph_slide(prs, title, img1_path, img2_path, top_color, text1, text2):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # BG
    f = slide.background.fill; f.solid(); f.fore_color.rgb = BG_C
    
    # Top bar
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(0.05))
    s.line.fill.background(); s.fill.solid(); s.fill.fore_color.rgb = top_color
    
    # Title
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(12), Inches(0.5))
    p = tb.text_frame.paragraphs[0]; p.text = title
    p.font.size = Pt(28); p.font.color.rgb = WHT_C; p.font.bold = True
    
    # Subtitle bar
    s2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(0.85), Inches(3), Inches(0.03))
    s2.line.fill.background(); s2.fill.solid(); s2.fill.fore_color.rgb = top_color

    # Insert images side by side
    # img1
    slide.shapes.add_picture(img1_path, Inches(0.5), Inches(1.8), height=Inches(4))
    # text 1
    t1 = slide.shapes.add_textbox(Inches(0.5), Inches(6.0), Inches(5.8), Inches(0.8))
    t1.text_frame.paragraphs[0].text = text1
    t1.text_frame.paragraphs[0].font.size = Pt(13); t1.text_frame.paragraphs[0].font.color.rgb = WHT_C
    t1.text_frame.word_wrap = True

    # img2
    slide.shapes.add_picture(img2_path, Inches(6.5), Inches(1.8), height=Inches(4))
    # text 2
    t2 = slide.shapes.add_textbox(Inches(6.5), Inches(6.0), Inches(5.8), Inches(0.8))
    t2.text_frame.paragraphs[0].text = text2
    t2.text_frame.paragraphs[0].font.size = Pt(13); t2.text_frame.paragraphs[0].font.color.rgb = WHT_C
    t2.text_frame.word_wrap = True
    
    # Bottom Bar
    s3 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.45), prs.slide_width, Inches(0.05))
    s3.line.fill.background(); s3.fill.solid(); s3.fill.fore_color.rgb = top_color

# Slide 6
add_graph_slide(
    prs, 
    "Data Visualized: Core Traffic Patterns", 
    "graphs/1_peak_hours.png", 
    "graphs/2_speed_seg.png",
    BLUE_C,
    "Insights \u2192 A huge spike in vehicle count exactly at 18:00 (6 PM) marks the clear evening rush hour. Deployment of adaptive signals should be prioritized between 16:00 and 20:00.",
    "Insights \u2192 Segment A1 exhibits the lowest network average speed. The infrastructure on this road fails to sustain flow capacity during high-demand bursts."
)

# Slide 7
add_graph_slide(
    prs, 
    "Data Visualized: Efficiency & Emissions", 
    "graphs/3_co2_cong.png", 
    "graphs/4_tei_rank.png",
    RGBColor(0xFF, 0x5C, 0x5C), # RED
    "Insights \u2192 Linear positive correlation between congestion states and CO2 emissions. Combating congestion directly translates to heavy pollution reduction.",
    "Insights \u2192 Traffic Efficiency Index (Speed/Density) reveals that C1 has the absolute worst operational efficiency, while D2 manages traffic the best."
)

# Slide 8
add_graph_slide(
    prs, 
    "Data Visualized: Future Simulations", 
    "graphs/6_ev_sim.png", 
    "graphs/5_glosa.png",
    ACCENT_C,
    "Insights \u2192 Simulating rapid EV adoption demonstrates that pushing the fleet to 50% EV immediately slashes network-wide CO2 levels by nearly 43%.",
    "Insights \u2192 Applying GLOSA to the worst bottlenecks (A1 & C1) visibly crushes CO2 spikes, preventing the stop-and-go shockwaves responsible for idle emissions."
)

# Move the 3 new slides to just BEFORE the CONCLUSION slide.
# The Conclusion slide is Currently at index 4 (5th slide).
# Our new slides are at indices 5, 6, 7.
slides = list(prs.slides._sldIdLst)
# Sequence: 0(Title), 1(Problem), 2(Analysis), 3(Solution), 4(Conclusion), 5(G1), 6(G2), 7(G3)
# Desired: 0, 1, 2, 3, 5, 6, 7, 4
prs.slides._sldIdLst.remove(slides[4])
prs.slides._sldIdLst.append(slides[4])

prs.save(pptx_out)

print("=" * 60)
print("  SUCCESS: GRAPHS INCLUDED IN NEW PPT!")
print("=" * 60)
print(f"  File name: Final_Presentation_With_Graphs.pptx")
print(f"  Total slides: 8")
print(f"  Location : {pptx_out}")
