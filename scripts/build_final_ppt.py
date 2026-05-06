import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import warnings
warnings.filterwarnings('ignore')

print("Building Executive Pitch PPT...")

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

# ---- COLORS ----
BG       = RGBColor(0x0A, 0x0F, 0x1A)   # deep navy
BG2      = RGBColor(0x10, 0x17, 0x28)
ACCENT   = RGBColor(0x00, 0xD2, 0xA8)   # bright teal
BLUE     = RGBColor(0x45, 0x7B, 0xFF)
RED      = RGBColor(0xFF, 0x5C, 0x5C)
YELLOW   = RGBColor(0xFF, 0xC1, 0x07)
GREEN    = RGBColor(0x00, 0xE6, 0x76)
PURPLE   = RGBColor(0xBB, 0x86, 0xFC)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
GREY     = RGBColor(0x7A, 0x83, 0x92)
LGREY    = RGBColor(0xBB, 0xBB, 0xBB)
CARD_BG  = RGBColor(0x14, 0x1C, 0x30)
DIVIDER  = RGBColor(0x1E, 0x29, 0x3E)

def set_bg(slide):
    f = slide.background.fill; f.solid(); f.fore_color.rgb = BG

def rect(slide, l, t, w, h, color=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.line.fill.background()
    if color: s.fill.solid(); s.fill.fore_color.rgb = color
    return s

def rrect(slide, l, t, w, h, color=None):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.line.fill.background()
    if color: s.fill.solid(); s.fill.fore_color.rgb = color
    return s

def txt(slide, l, t, w, h, text, sz=14, color=WHITE, bold=False, italic=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz); p.font.color.rgb = color; p.font.bold = bold
    p.font.italic = italic; p.font.name = "Calibri"; p.alignment = align
    return tb

def multi(slide, l, t, w, h, lines):
    """lines = [(text, size, color, bold, align), ...]"""
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    for i, (text, sz, c, b, a) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text; p.font.size = Pt(sz); p.font.color.rgb = c
        p.font.bold = b; p.font.name = "Calibri"; p.alignment = a
        p.space_before = Pt(4); p.space_after = Pt(4)

def big_number(slide, l, t, w, number, label, accent_color):
    """Big stat card"""
    rrect(slide, l, t, w, Inches(1.5), CARD_BG)
    rect(slide, l, t, w, Inches(0.05), accent_color)
    txt(slide, l + Inches(0.2), t + Inches(0.15), w - Inches(0.4), Inches(0.7),
        number, sz=32, color=accent_color, bold=True, align=PP_ALIGN.CENTER)
    txt(slide, l + Inches(0.2), t + Inches(0.95), w - Inches(0.4), Inches(0.4),
        label, sz=10, color=GREY, align=PP_ALIGN.CENTER)

W = prs.slide_width

# ============================================================
# SLIDE 1: TITLE
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
rect(s, Inches(0), Inches(0), W, Inches(0.05), ACCENT)

# Centered title block
txt(s, Inches(1), Inches(1.8), Inches(11.3), Inches(1),
    "Smart Highway Traffic Analysis",
    sz=42, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

txt(s, Inches(1), Inches(2.9), Inches(11.3), Inches(0.6),
    "Data-Driven Insights for Emission Reduction & Congestion Management",
    sz=18, color=ACCENT, italic=True, align=PP_ALIGN.CENTER)

# Thin divider line
rect(s, Inches(5), Inches(3.7), Inches(3.3), Inches(0.02), ACCENT)

txt(s, Inches(1), Inches(4.0), Inches(11.3), Inches(0.4),
    "Highway Dataset  |  4,012 Records  |  8 Road Segments  |  Jan-Mar 2026",
    sz=12, color=GREY, align=PP_ALIGN.CENTER)

# Bottom stat strip
rect(s, Inches(0), Inches(5.8), W, Inches(0.02), DIVIDER)
stats = [
    ("4,012", "Clean Records", ACCENT),
    ("108", "Duplicates Removed", RED),
    ("24", "Features", BLUE),
    ("8", "Segments", PURPLE),
    ("13", "KPIs Computed", YELLOW),
]
for i, (val, label, clr) in enumerate(stats):
    x = Inches(0.8 + i * 2.5)
    txt(s, x, Inches(6.0), Inches(2), Inches(0.5),
        val, sz=28, color=clr, bold=True, align=PP_ALIGN.CENTER)
    txt(s, x, Inches(6.5), Inches(2), Inches(0.3),
        label, sz=9, color=GREY, align=PP_ALIGN.CENTER)

rect(s, Inches(0), Inches(7.45), W, Inches(0.05), ACCENT)


# ============================================================
# SLIDE 2: THE PROBLEM — WHAT THE DATA TELLS US
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
rect(s, Inches(0), Inches(0), W, Inches(0.05), ACCENT)

txt(s, Inches(0.6), Inches(0.3), Inches(12), Inches(0.5),
    "The Problem: What the Data Reveals",
    sz=28, color=WHITE, bold=True)
rect(s, Inches(0.6), Inches(0.85), Inches(3), Inches(0.03), ACCENT)

# 5 Big stat cards — top row
big_number(s, Inches(0.4),  Inches(1.2), Inches(2.4), "59.72", "Avg Speed (kmph)", BLUE)
big_number(s, Inches(3.05), Inches(1.2), Inches(2.4), "249.08", "Avg CO\u2082 Emission", RED)
big_number(s, Inches(5.7),  Inches(1.2), Inches(2.4), "18:00", "Peak Traffic Hour", YELLOW)
big_number(s, Inches(8.35), Inches(1.2), Inches(2.4), "A1", "Slowest Segment", RED)
big_number(s, Inches(11.0), Inches(1.2), Inches(2.4), "D1", "Most Congested", PURPLE)

# Key findings in 3 columns
col_w = Inches(3.9)

# Column 1: Congestion
rrect(s, Inches(0.4), Inches(3.0), col_w, Inches(3.8), CARD_BG)
rect(s, Inches(0.4), Inches(3.0), col_w, Inches(0.04), RED)
multi(s, Inches(0.7), Inches(3.2), Inches(3.3), Inches(3.5), [
    ("Congestion = Pollution", 16, RED, True, PP_ALIGN.LEFT),
    ("", 8, WHITE, False, PP_ALIGN.LEFT),
    ("High congestion areas produce", 12, LGREY, False, PP_ALIGN.LEFT),
    ("18% more CO\u2082 than low congestion", 14, WHITE, True, PP_ALIGN.LEFT),
    ("", 8, WHITE, False, PP_ALIGN.LEFT),
    ("High:    266.86 avg CO\u2082", 11, RED, False, PP_ALIGN.LEFT),
    ("Medium: 252.34 avg CO\u2082", 11, YELLOW, False, PP_ALIGN.LEFT),
    ("Low:     226.01 avg CO\u2082", 11, GREEN, False, PP_ALIGN.LEFT),
    ("", 8, WHITE, False, PP_ALIGN.LEFT),
    ("Direct link: reduce congestion", 10, GREY, False, PP_ALIGN.LEFT),
    ("= reduce pollution", 10, GREY, False, PP_ALIGN.LEFT),
])

# Column 2: Worst Segments
rrect(s, Inches(4.6), Inches(3.0), col_w, Inches(3.8), CARD_BG)
rect(s, Inches(4.6), Inches(3.0), col_w, Inches(0.04), BLUE)
multi(s, Inches(4.9), Inches(3.2), Inches(3.3), Inches(3.5), [
    ("Worst Performing Roads", 16, BLUE, True, PP_ALIGN.LEFT),
    ("", 8, WHITE, False, PP_ALIGN.LEFT),
    ("Segment A1: Slowest (57.46 kmph)", 12, RED, False, PP_ALIGN.LEFT),
    ("Segment D1: Most Dense (71.54)", 12, RED, False, PP_ALIGN.LEFT),
    ("Segment C1: Worst Efficiency", 12, RED, False, PP_ALIGN.LEFT),
    ("", 8, WHITE, False, PP_ALIGN.LEFT),
    ("Top Emission Hotspot:", 12, LGREY, False, PP_ALIGN.LEFT),
    ("Segment A1 appears TWICE", 14, YELLOW, True, PP_ALIGN.LEFT),
    ("in the top 5 CO\u2082 hotspots", 12, YELLOW, False, PP_ALIGN.LEFT),
    ("", 8, WHITE, False, PP_ALIGN.LEFT),
    ("A1 needs urgent intervention", 10, GREY, False, PP_ALIGN.LEFT),
])

# Column 3: Traffic Patterns
rrect(s, Inches(8.8), Inches(3.0), col_w, Inches(3.8), CARD_BG)
rect(s, Inches(8.8), Inches(3.0), col_w, Inches(0.04), ACCENT)
multi(s, Inches(9.1), Inches(3.2), Inches(3.3), Inches(3.5), [
    ("Traffic Patterns", 16, ACCENT, True, PP_ALIGN.LEFT),
    ("", 8, WHITE, False, PP_ALIGN.LEFT),
    ("Peak hour: 6 PM (327,645 vehicles)", 12, LGREY, False, PP_ALIGN.LEFT),
    ("Evening rush dominates (6-10 PM)", 12, LGREY, False, PP_ALIGN.LEFT),
    ("", 8, WHITE, False, PP_ALIGN.LEFT),
    ("Busiest day: Saturday", 12, LGREY, False, PP_ALIGN.LEFT),
    ("Quietest day: Wednesday", 12, LGREY, False, PP_ALIGN.LEFT),
    ("", 8, WHITE, False, PP_ALIGN.LEFT),
    ("Surprising finding:", 12, YELLOW, True, PP_ALIGN.LEFT),
    ("Rainfall has NO significant", 12, WHITE, False, PP_ALIGN.LEFT),
    ("impact on speed (<2 kmph diff)", 12, WHITE, False, PP_ALIGN.LEFT),
])

rect(s, Inches(0), Inches(7.45), W, Inches(0.05), ACCENT)


# ============================================================
# SLIDE 3: THE NUMBERS — DATA TELLS THE TRUTH
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
rect(s, Inches(0), Inches(0), W, Inches(0.05), BLUE)

txt(s, Inches(0.6), Inches(0.3), Inches(12), Inches(0.5),
    "Deep Analysis: Efficiency, EV Impact & Gore Points",
    sz=28, color=WHITE, bold=True)
rect(s, Inches(0.6), Inches(0.85), Inches(3), Inches(0.03), BLUE)

# TEI Ranking — left section
rrect(s, Inches(0.4), Inches(1.2), Inches(4.0), Inches(5.6), CARD_BG)
rect(s, Inches(0.4), Inches(1.2), Inches(4.0), Inches(0.04), PURPLE)
multi(s, Inches(0.7), Inches(1.4), Inches(3.4), Inches(5.2), [
    ("Traffic Efficiency Index", 16, PURPLE, True, PP_ALIGN.LEFT),
    ("TEI = Speed / Density", 10, GREY, False, PP_ALIGN.LEFT),
    ("Higher = better road performance", 10, GREY, False, PP_ALIGN.LEFT),
    ("", 8, WHITE, False, PP_ALIGN.LEFT),
    ("#1  D2    TEI = 1.390  (BEST)", 12, GREEN, True, PP_ALIGN.LEFT),
    ("#2  A2    TEI = 1.382", 12, LGREY, False, PP_ALIGN.LEFT),
    ("#3  B1    TEI = 1.378", 12, LGREY, False, PP_ALIGN.LEFT),
    ("#4  C2    TEI = 1.363", 12, LGREY, False, PP_ALIGN.LEFT),
    ("#5  B2    TEI = 1.340", 12, LGREY, False, PP_ALIGN.LEFT),
    ("#6  D1    TEI = 1.331", 12, LGREY, False, PP_ALIGN.LEFT),
    ("#7  A1    TEI = 1.279", 12, YELLOW, False, PP_ALIGN.LEFT),
    ("#8  C1    TEI = 1.240  (WORST)", 12, RED, True, PP_ALIGN.LEFT),
    ("", 8, WHITE, False, PP_ALIGN.LEFT),
    ("Study D2 config as model", 10, ACCENT, False, PP_ALIGN.LEFT),
    ("Redesign C1 infrastructure", 10, RED, False, PP_ALIGN.LEFT),
])

# EV Simulation — center section
rrect(s, Inches(4.7), Inches(1.2), Inches(4.0), Inches(5.6), CARD_BG)
rect(s, Inches(4.7), Inches(1.2), Inches(4.0), Inches(0.04), GREEN)
multi(s, Inches(5.0), Inches(1.4), Inches(3.4), Inches(5.2), [
    ("EV Adoption Simulation", 16, GREEN, True, PP_ALIGN.LEFT),
    ("What if we increase EVs?", 10, GREY, False, PP_ALIGN.LEFT),
    ("", 8, WHITE, False, PP_ALIGN.LEFT),
    ("Current fleet: 33.9% EV", 12, LGREY, False, PP_ALIGN.LEFT),
    ("Current avg CO\u2082: 249.08", 12, LGREY, False, PP_ALIGN.LEFT),
    ("", 8, WHITE, False, PP_ALIGN.LEFT),
    ("25% EV  \u2192  21.3% CO\u2082 reduction", 13, LGREY, False, PP_ALIGN.LEFT),
    ("50% EV  \u2192  42.5% CO\u2082 reduction", 13, GREEN, True, PP_ALIGN.LEFT),
    ("75% EV  \u2192  63.7% CO\u2082 reduction", 13, LGREY, False, PP_ALIGN.LEFT),
    ("100% EV \u2192  85.0% CO\u2082 reduction", 13, ACCENT, True, PP_ALIGN.LEFT),
    ("", 8, WHITE, False, PP_ALIGN.LEFT),
    ("Key takeaway:", 12, YELLOW, True, PP_ALIGN.LEFT),
    ("Just reaching 50% EV adoption", 12, WHITE, False, PP_ALIGN.LEFT),
    ("cuts nearly HALF of all CO\u2082", 12, WHITE, False, PP_ALIGN.LEFT),
    ("", 8, WHITE, False, PP_ALIGN.LEFT),
    ("Build charging infra NOW", 10, GREEN, False, PP_ALIGN.LEFT),
])

# Gore Point — right section
rrect(s, Inches(9.0), Inches(1.2), Inches(4.0), Inches(2.5), CARD_BG)
rect(s, Inches(9.0), Inches(1.2), Inches(4.0), Inches(0.04), YELLOW)
multi(s, Inches(9.3), Inches(1.4), Inches(3.4), Inches(2.2), [
    ("Gore Point Impact", 16, YELLOW, True, PP_ALIGN.LEFT),
    ("Merge/diverge zones on highway", 10, GREY, False, PP_ALIGN.LEFT),
    ("", 8, WHITE, False, PP_ALIGN.LEFT),
    ("           Non-Gore    Gore", 10, GREY, False, PP_ALIGN.LEFT),
    ("Speed:    58.4        61.4  kmph", 12, LGREY, False, PP_ALIGN.LEFT),
    ("CO\u2082:      244.1       255.6", 12, LGREY, False, PP_ALIGN.LEFT),
    ("Density: 64.6        69.2", 12, LGREY, False, PP_ALIGN.LEFT),
])

# Bonus insights — bottom right
rrect(s, Inches(9.0), Inches(4.0), Inches(4.0), Inches(2.8), CARD_BG)
rect(s, Inches(9.0), Inches(4.0), Inches(4.0), Inches(0.04), ACCENT)
multi(s, Inches(9.3), Inches(4.2), Inches(3.4), Inches(2.5), [
    ("New Findings for Policy Board", 16, ACCENT, True, PP_ALIGN.LEFT),
    ("", 6, WHITE, False, PP_ALIGN.LEFT),
    ("1. Saturday is busiest day", 11, WHITE, False, PP_ALIGN.LEFT),
    ("   \u2192 Schedule roadwork on Wednesday", 10, GREY, False, PP_ALIGN.LEFT),
    ("", 6, WHITE, False, PP_ALIGN.LEFT),
    ("2. Rain doesn't slow traffic", 11, WHITE, False, PP_ALIGN.LEFT),
    ("   \u2192 Density is the real bottleneck", 10, GREY, False, PP_ALIGN.LEFT),
    ("", 6, WHITE, False, PP_ALIGN.LEFT),
    ("3. A1 appears in worst lists 3 times", 11, RED, False, PP_ALIGN.LEFT),
    ("   \u2192 Priority #1 for intervention", 10, RED, False, PP_ALIGN.LEFT),
])

rect(s, Inches(0), Inches(7.45), W, Inches(0.05), BLUE)


# ============================================================
# SLIDE 4: THE SOLUTION — GLOSA-BHARAT
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
rect(s, Inches(0), Inches(0), W, Inches(0.05), GREEN)

txt(s, Inches(0.6), Inches(0.3), Inches(12), Inches(0.5),
    "The Solution: GLOSA-BHARAT Deployment Simulation",
    sz=28, color=WHITE, bold=True)
rect(s, Inches(0.6), Inches(0.85), Inches(3), Inches(0.03), GREEN)

# 4 impact cards
big_number(s, Inches(0.4),  Inches(1.2), Inches(3.0), "35,416", "CO\u2082 Units Saved", GREEN)
big_number(s, Inches(3.65), Inches(1.2), Inches(3.0), "+2.3%", "Speed Improvement", ACCENT)
big_number(s, Inches(6.9),  Inches(1.2), Inches(3.0), "A1 & C1", "Segments Targeted", RED)
big_number(s, Inches(10.15),Inches(1.2), Inches(3.0), "532", "Rows Affected", BLUE)

# What is GLOSA
rrect(s, Inches(0.4), Inches(3.0), Inches(4.0), Inches(3.8), CARD_BG)
rect(s, Inches(0.4), Inches(3.0), Inches(4.0), Inches(0.04), GREEN)
multi(s, Inches(0.7), Inches(3.2), Inches(3.4), Inches(3.5), [
    ("What is GLOSA?", 18, GREEN, True, PP_ALIGN.LEFT),
    ("", 6, WHITE, False, PP_ALIGN.LEFT),
    ("Green Light Optimized", 13, WHITE, False, PP_ALIGN.LEFT),
    ("Speed Advisory", 13, WHITE, False, PP_ALIGN.LEFT),
    ("", 6, WHITE, False, PP_ALIGN.LEFT),
    ("A V2X system that tells drivers", 11, LGREY, False, PP_ALIGN.LEFT),
    ("the optimal speed to maintain,", 11, LGREY, False, PP_ALIGN.LEFT),
    ("preventing stop-and-go traffic.", 11, LGREY, False, PP_ALIGN.LEFT),
    ("", 6, WHITE, False, PP_ALIGN.LEFT),
    ("No sudden braking", 12, ACCENT, True, PP_ALIGN.LEFT),
    ("= No shockwave jams", 12, ACCENT, True, PP_ALIGN.LEFT),
    ("= No idle CO\u2082 emissions", 12, ACCENT, True, PP_ALIGN.LEFT),
])

# Simulation before/after
rrect(s, Inches(4.7), Inches(3.0), Inches(4.0), Inches(3.8), CARD_BG)
rect(s, Inches(4.7), Inches(3.0), Inches(4.0), Inches(0.04), BLUE)
multi(s, Inches(5.0), Inches(3.2), Inches(3.4), Inches(3.5), [
    ("Before vs After GLOSA", 18, BLUE, True, PP_ALIGN.LEFT),
    ("", 6, WHITE, False, PP_ALIGN.LEFT),
    ("SEGMENT A1:", 14, WHITE, True, PP_ALIGN.LEFT),
    ("Speed: 57.5 \u2192 63.3 kmph (+10.2%)", 12, GREEN, False, PP_ALIGN.LEFT),
    ("CO\u2082:   244.5 \u2192 212.3 (-13.2%)", 12, GREEN, False, PP_ALIGN.LEFT),
    ("", 6, WHITE, False, PP_ALIGN.LEFT),
    ("SEGMENT C1:", 14, WHITE, True, PP_ALIGN.LEFT),
    ("Speed: 57.6 \u2192 64.1 kmph (+11.3%)", 12, GREEN, False, PP_ALIGN.LEFT),
    ("CO\u2082:   252.0 \u2192 215.8 (-14.4%)", 12, GREEN, False, PP_ALIGN.LEFT),
    ("", 6, WHITE, False, PP_ALIGN.LEFT),
    ("Just 2 segments targeted,", 11, YELLOW, True, PP_ALIGN.LEFT),
    ("3.5% total network CO\u2082 drop", 11, YELLOW, True, PP_ALIGN.LEFT),
])

# Why it matters
rrect(s, Inches(9.0), Inches(3.0), Inches(4.0), Inches(3.8), CARD_BG)
rect(s, Inches(9.0), Inches(3.0), Inches(4.0), Inches(0.04), YELLOW)
multi(s, Inches(9.3), Inches(3.2), Inches(3.4), Inches(3.5), [
    ("Policy Recommendation", 18, YELLOW, True, PP_ALIGN.LEFT),
    ("", 6, WHITE, False, PP_ALIGN.LEFT),
    ("IMMEDIATE:", 12, YELLOW, True, PP_ALIGN.LEFT),
    ("Deploy GLOSA on A1 & C1", 11, LGREY, False, PP_ALIGN.LEFT),
    ("Adaptive signals at 6 PM peak", 11, LGREY, False, PP_ALIGN.LEFT),
    ("", 6, WHITE, False, PP_ALIGN.LEFT),
    ("SHORT-TERM:", 12, ACCENT, True, PP_ALIGN.LEFT),
    ("EV incentives (target 50%)", 11, LGREY, False, PP_ALIGN.LEFT),
    ("Restructure C1 (worst TEI)", 11, LGREY, False, PP_ALIGN.LEFT),
    ("", 6, WHITE, False, PP_ALIGN.LEFT),
    ("LONG-TERM:", 12, PURPLE, True, PP_ALIGN.LEFT),
    ("GLOSA on all 8 segments", 11, LGREY, False, PP_ALIGN.LEFT),
    ("Weather-adaptive speed advisories", 11, LGREY, False, PP_ALIGN.LEFT),
])

rect(s, Inches(0), Inches(7.45), W, Inches(0.05), GREEN)


# ============================================================
# SLIDE 5: CONCLUSION
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
rect(s, Inches(0), Inches(0), W, Inches(0.05), ACCENT)

txt(s, Inches(1), Inches(0.8), Inches(11.3), Inches(0.8),
    "We didn't just analyse the data.",
    sz=32, color=GREY, italic=True, align=PP_ALIGN.CENTER)

txt(s, Inches(1), Inches(1.7), Inches(11.3), Inches(0.8),
    "We proved the solution works.",
    sz=36, color=GREEN, bold=True, align=PP_ALIGN.CENTER)

rect(s, Inches(5), Inches(2.7), Inches(3.3), Inches(0.02), ACCENT)

# Summary stats row
summary = [
    ("6/8", "Core Analyses", BLUE),
    ("3/5", "Advanced Analyses", PURPLE),
    ("13", "KPIs Computed", YELLOW),
    ("35,416", "CO\u2082 Saved by GLOSA", GREEN),
    ("+42.5%", "CO\u2082 Cut at 50% EV", ACCENT),
]

for i, (val, label, clr) in enumerate(summary):
    x = Inches(0.5 + i * 2.55)
    big_number(s, x, Inches(3.1), Inches(2.3), val, label, clr)

# What we delivered
rrect(s, Inches(0.5), Inches(5.0), Inches(6.0), Inches(1.8), CARD_BG)
rect(s, Inches(0.5), Inches(5.0), Inches(6.0), Inches(0.04), BLUE)
multi(s, Inches(0.8), Inches(5.15), Inches(5.4), Inches(1.6), [
    ("What We Delivered", 16, BLUE, True, PP_ALIGN.LEFT),
    ("", 4, WHITE, False, PP_ALIGN.LEFT),
    ("\u2713  Data cleaned: 108 duplicates removed, all blanks imputed", 11, LGREY, False, PP_ALIGN.LEFT),
    ("\u2713  Features: Hour, TEI, Rain Category engineered", 11, LGREY, False, PP_ALIGN.LEFT),
    ("\u2713  6 Core + 3 Advanced + 2 Bonus analyses completed", 11, LGREY, False, PP_ALIGN.LEFT),
    ("\u2713  Excel Dashboard: 10 charts, 5 sheets, 13 KPIs", 11, LGREY, False, PP_ALIGN.LEFT),
    ("\u2713  GLOSA-BHARAT simulation with proven CO\u2082 savings", 11, GREEN, False, PP_ALIGN.LEFT),
])

# Call to action
rrect(s, Inches(6.8), Inches(5.0), Inches(6.0), Inches(1.8), CARD_BG)
rect(s, Inches(6.8), Inches(5.0), Inches(6.0), Inches(0.04), GREEN)
multi(s, Inches(7.1), Inches(5.15), Inches(5.4), Inches(1.6), [
    ("The Opportunity", 16, GREEN, True, PP_ALIGN.LEFT),
    ("", 4, WHITE, False, PP_ALIGN.LEFT),
    ("If GLOSA is deployed on just 2 worst segments,", 12, WHITE, False, PP_ALIGN.LEFT),
    ("we save 35,416 CO\u2082 units (3.5% total).", 12, WHITE, False, PP_ALIGN.LEFT),
    ("", 4, WHITE, False, PP_ALIGN.LEFT),
    ("Scale to all 8 segments:", 12, YELLOW, True, PP_ALIGN.LEFT),
    ("Estimated 10-15% total CO\u2082 reduction", 14, GREEN, True, PP_ALIGN.LEFT),
    ("without building a single new road.", 12, ACCENT, False, PP_ALIGN.LEFT),
])

rect(s, Inches(0), Inches(7.45), W, Inches(0.05), ACCENT)

# ============================================================
# SAVE
# ============================================================
out = r'c:\Users\name\Desktop\PROJECTS\GLOSA-BHARAT-main\Final_Presentation.pptx'
prs.save(out)
print(f"  SAVED: Final_Presentation.pptx")
print(f"  5 slides | Executive Pitch style")
print(f"  Ready to present!")
